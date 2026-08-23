from __future__ import annotations

import csv
import re
from dataclasses import dataclass
from io import BytesIO, StringIO

from docx import Document
from docx.enum.text import WD_BREAK
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt

from .schemas import ComparativeMatrixRead, ExtractionParameter


class ExportUnavailableError(Exception):
    """Raised when a project has no reviewed, ready result to export."""


@dataclass(frozen=True)
class ExportArtifact:
    content: bytes
    filename: str
    media_type: str


PARAMETER_LABELS = {
    ExtractionParameter.RESEARCH_PROBLEM: "Research problem",
    ExtractionParameter.RESEARCH_OBJECTIVE: "Research objective",
    ExtractionParameter.RESEARCH_QUESTION: "Research question",
    ExtractionParameter.METHODOLOGY: "Methodology",
    ExtractionParameter.DATASET_SAMPLE: "Dataset / sample",
    ExtractionParameter.VARIABLES_CONCEPTS: "Variables / concepts",
    ExtractionParameter.RESULTS_FINDINGS: "Results / findings",
    ExtractionParameter.CONTRIBUTION: "Contribution",
    ExtractionParameter.LIMITATIONS: "Limitations",
    ExtractionParameter.FUTURE_WORK: "Future work",
    ExtractionParameter.KEY_CLAIMS: "Key claims",
}


def build_project_export(matrix: ComparativeMatrixRead, export_format: str) -> ExportArtifact:
    if not matrix.papers:
        raise ExportUnavailableError(
            "Belum ada paper siap yang dapat diekspor. Selesaikan review paper dahulu."
        )
    builders = {
        "docx": (_build_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        "xlsx": (_build_xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        "csv": (_build_csv, "text/csv; charset=utf-8"),
        "pptx": (_build_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    }
    if export_format not in builders:
        raise ValueError("Format ekspor harus docx, xlsx, csv, atau pptx")
    builder, media_type = builders[export_format]
    slug = re.sub(r"[^a-z0-9]+", "-", matrix.project_title.lower()).strip("-")
    filename = f"modellyng-{slug or 'project'}.{export_format}"
    return ExportArtifact(builder(matrix), filename, media_type)


def _clean(value: object, limit: int = 32_000) -> str:
    text = "" if value is None else str(value)
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)
    return text[:limit]


def _records(matrix: ComparativeMatrixRead):
    paper_by_id = {paper.id: paper for paper in matrix.papers}
    for row in matrix.rows:
        for cell in row.cells:
            paper = paper_by_id.get(cell.paper_id)
            if paper is None:
                continue
            value = cell.final_value if cell.final_value is not None else cell.ai_value
            evidence = "\n".join(
                f'Page {item.page_number}: "{item.quote}" [block {item.block_id}]'
                for item in cell.evidence
            )
            yield {
                "paper_id": str(paper.id),
                "paper": paper.title,
                "filename": paper.original_filename,
                "parameter": PARAMETER_LABELS[row.parameter],
                "parameter_key": row.parameter.value,
                "reviewed_value": value,
                "original_ai_value": cell.ai_value,
                "status": cell.status.value,
                "confidence": "" if cell.confidence is None else f"{cell.confidence:.0%}",
                "evidence": evidence,
            }


def _build_csv(matrix: ComparativeMatrixRead) -> bytes:
    output = StringIO(newline="")
    columns = [
        "paper_id", "paper", "filename", "parameter", "parameter_key",
        "reviewed_value", "original_ai_value", "status", "confidence", "evidence",
    ]
    writer = csv.DictWriter(output, fieldnames=columns)
    writer.writeheader()
    for record in _records(matrix):
        writer.writerow({key: _clean(value) for key, value in record.items()})
    return ("\ufeff" + output.getvalue()).encode("utf-8")


def _build_xlsx(matrix: ComparativeMatrixRead) -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Reviewed results"
    headers = [
        "Paper ID", "Paper", "Source file", "Parameter", "Parameter key",
        "Reviewed value", "Original AI value", "Status", "Confidence", "Evidence",
    ]
    sheet.append(headers)
    for record in _records(matrix):
        sheet.append([_clean(value) for value in record.values()])
    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="5147E5")
    sheet.freeze_panes = "A2"
    sheet.auto_filter.ref = sheet.dimensions
    widths = [38, 28, 26, 24, 24, 55, 55, 18, 14, 70]
    for index, width in enumerate(widths, 1):
        sheet.column_dimensions[chr(64 + index)].width = width
    for row in sheet.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
    notes = workbook.create_sheet("Read me")
    notes.append(["Modellyng evidence-preserving export"])
    notes.append(["Project", _clean(matrix.project_title)])
    notes.append(["Scope", "Ready papers and their latest reviewed components only."])
    notes.append(["Evidence", "Quotes retain page and private paper block identifiers."])
    notes.append(["Human review", "Reviewed value is final_value when edited; otherwise the accepted AI value."])
    notes["A1"].font = Font(bold=True, color="5147E5", size=14)
    stream = BytesIO()
    workbook.save(stream)
    return stream.getvalue()


def _build_docx(matrix: ComparativeMatrixRead) -> bytes:
    document = Document()
    document.add_heading(matrix.project_title, 0)
    document.add_paragraph(
        "Modellyng reviewed result export. Only ready papers are included. "
        "Evidence quotes retain their source page, block ID, and private paper ID."
    )
    records = list(_records(matrix))
    for paper_index, paper in enumerate(matrix.papers):
        document.add_heading(paper.title, level=1)
        document.add_paragraph(f"Source file: {paper.original_filename}\nPaper ID: {paper.id}")
        for record in (item for item in records if item["paper_id"] == str(paper.id)):
            document.add_heading(_clean(record["parameter"]), level=2)
            document.add_paragraph(_clean(record["reviewed_value"]))
            metadata = document.add_paragraph()
            metadata.add_run(
                f"Status: {record['status']} · Confidence: {record['confidence'] or 'n/a'}"
            ).bold = True
            if record["reviewed_value"] != record["original_ai_value"]:
                document.add_paragraph(
                    f"Original AI value: {_clean(record['original_ai_value'])}"
                )
            if record["evidence"]:
                document.add_paragraph("Evidence", style="Intense Quote")
                for evidence_line in str(record["evidence"]).splitlines():
                    document.add_paragraph(_clean(evidence_line), style="List Bullet")
        if paper_index < len(matrix.papers) - 1:
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def _build_pptx(matrix: ComparativeMatrixRead) -> bytes:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = _clean(matrix.project_title, 200)
    title_slide.placeholders[1].text = (
        "Reviewed academic results\nEvidence-preserving Modellyng export"
    )
    records = list(_records(matrix))
    for paper in matrix.papers:
        paper_records = [item for item in records if item["paper_id"] == str(paper.id)]
        for offset in range(0, len(paper_records), 3):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = _clean(paper.title, 150)
            subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.8), Inches(0.35))
            subtitle.text_frame.text = f"{paper.original_filename} · Paper ID {paper.id}"
            top = 1.65
            for record in paper_records[offset : offset + 3]:
                box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.8), Inches(1.55))
                frame = box.text_frame
                frame.word_wrap = True
                paragraph = frame.paragraphs[0]
                paragraph.text = _clean(record["parameter"], 80)
                paragraph.font.bold = True
                paragraph.font.size = Pt(18)
                value = frame.add_paragraph()
                value.text = _clean(record["reviewed_value"], 600)
                value.font.size = Pt(13)
                evidence = frame.add_paragraph()
                evidence.text = _clean(
                    f"{record['status']} · confidence {record['confidence'] or 'n/a'} · "
                    f"{record['evidence'] or 'No verified quote'}",
                    500,
                )
                evidence.font.size = Pt(10)
                evidence.font.italic = True
                top += 1.75
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()
